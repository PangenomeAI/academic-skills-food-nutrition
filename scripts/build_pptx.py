#!/usr/bin/env python3
"""Build a fully EDITABLE, well-designed PowerPoint (.pptx) from a deck-spec JSON.

Everything on every slide is a native, editable object: titles/body are real text
frames, tables are native PowerPoint tables, figures are movable picture objects, and
the visual design (colour panels, accent stripes, section dividers, footers, slide
numbers) is built from **editable shapes** — nothing is flattened to a background image.

Themes live in `food-ppt/templates/themes.json`; pick one with the deck-spec `"theme"`
field (default: the file's `default`). Requires python-pptx. The spec validator is
stdlib-only so `--selftest` runs without python-pptx.

Deck-spec JSON:
{
  "title": "...", "subtitle": "...", "authors": "...", "date": "...",
  "theme": "sage",                              // optional; a themes.json key
  "footer": "Short deck label",                 // optional; shown in the footer
  "slides": [
    {"layout": "title",   "title": "...", "subtitle": "...", "authors": "...", "notes": "..."},
    {"layout": "section", "title": "...", "notes": "..."},
    {"layout": "bullets", "title": "...", "bullets": ["a", ["sub"], "b"], "notes": "..."},
    {"layout": "two_column", "title": "...", "left": ["..."], "right": ["..."], "notes": "..."},
    {"layout": "figure",  "title": "...", "image": "fig.png", "caption": "...", "notes": "..."},
    {"layout": "table",   "title": "...", "table": [["H1","H2"], ["a","b"]], "notes": "..."},
    {"layout": "metric",  "title": "...", "value": "42%", "label": "reduction ...", "support": "vs control (p<0.05)", "notes": "..."},
    {"layout": "cards",   "title": "...", "cards": [{"heading": "n = 1,240", "text": "..."}], "notes": "..."},
    {"layout": "references", "title": "References", "items": ["..."], "notes": "..."}
  ]
}

Usage:
  build_pptx.py deck.json --out deck.pptx [--theme sage]
  build_pptx.py --selftest
"""
import json
import math
import os
import sys
import zipfile

LAYOUTS = {"title", "section", "bullets", "two_column", "figure", "table",
           "references", "flow", "result", "metric", "cards"}

_FALLBACK_THEME = {
    "bg": "FFFFFF", "surface": "EEF2F6", "band": "F5F8FB", "ink": "18212B",
    "muted": "5B6673", "accent": "1F3A5F", "accent2": "3E7CB1", "onaccent": "FFFFFF",
    "font_head": "Calibri", "font_body": "Calibri",
}


def themes_path():
    root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(root, "food-ppt", "templates", "themes.json")


def load_theme(name=None):
    try:
        data = json.load(open(themes_path(), encoding="utf-8"))
        themes = data.get("themes", {})
        key = name or data.get("default")
        return {**_FALLBACK_THEME, **themes.get(key, themes.get(data.get("default"), {}))}
    except (OSError, ValueError):
        return dict(_FALLBACK_THEME)


def validate_spec(spec):
    """Return a list of problems (stdlib only). Empty list == valid."""
    if not isinstance(spec, dict):
        return ["spec must be a JSON object"]
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        return ["spec.slides must be a non-empty list"]
    problems = []
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            problems.append(f"slide {i}: not an object"); continue
        lay = s.get("layout")
        if lay not in LAYOUTS:
            problems.append(f"slide {i}: layout '{lay}' not in {sorted(LAYOUTS)}")
        if lay == "bullets" and not s.get("bullets"):
            problems.append(f"slide {i} (bullets): missing 'bullets'")
        if lay == "figure" and not s.get("image"):
            problems.append(f"slide {i} (figure): missing 'image'")
        if lay == "table":
            t = s.get("table")
            if not (isinstance(t, list) and t and all(isinstance(r, list) for r in t)):
                problems.append(f"slide {i} (table): 'table' must be a list of rows")
        if lay == "references" and not s.get("items"):
            problems.append(f"slide {i} (references): missing 'items'")
        if lay == "flow" and not (s.get("steps") or s.get("lanes")):
            problems.append(f"slide {i} (flow): needs 'steps' (list) or 'lanes' (labelled rows)")
        if lay == "result":
            if not (s.get("image") or s.get("table")):
                problems.append(f"slide {i} (result): needs an 'image' or a 'table'")
            if not s.get("key_findings"):
                problems.append(f"slide {i} (result): needs 'key_findings' (list)")
        if lay == "metric" and not s.get("value"):
            problems.append(f"slide {i} (metric): needs a 'value' (the big number)")
        if lay == "cards":
            cards = s.get("cards")
            if not (isinstance(cards, list) and cards):
                problems.append(f"slide {i} (cards): needs 'cards' (list of {{heading,text}})")
    return problems


def build(spec, out_path, theme_name=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    th = load_theme(theme_name or spec.get("theme"))
    C = {k: RGBColor.from_string(th[k]) for k in
         ("bg", "surface", "band", "ink", "muted", "accent", "accent2", "onaccent")}
    FH, FB = th["font_head"], th["font_body"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    def rect(sl, l, t, w, h, fill, line=None):
        sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sp.shadow.inherit = False
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(0.75)
        return sp

    def tb(sl, l, t, w, h, anchor=None):
        box = sl.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        return tf

    def para(p, text, size, font=FB, bold=False, color=None, align=None, space_after=6):
        p.text = text
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
            r.font.color.rgb = color if color is not None else C["ink"]
        if align is not None:
            p.alignment = align
        p.space_after = Pt(space_after)
        return p

    def track(p, pts):
        """Letter-spacing (tracking) in points; negative = tight. Editable rPr attr."""
        for r in p.runs:
            r._r.get_or_add_rPr().set("spc", str(int(pts * 100)))
        return p

    def line_spacing(p, mult):
        p.line_spacing = mult
        return p

    def bg(sl, color=None):
        rect(sl, 0, 0, SW, SH, color or C["bg"])

    def footer(sl, n):
        rect(sl, Inches(0.7), Inches(7.02), Inches(11.93), Pt(0.9), C["surface"])
        ft = tb(sl, Inches(0.7), Inches(7.06), Inches(9), Inches(0.35))
        para(ft.paragraphs[0], spec.get("footer", spec.get("title", "")), 9,
             color=C["muted"], space_after=0)
        pn = tb(sl, Inches(11.9), Inches(7.06), Inches(0.7), Inches(0.35))
        para(pn.paragraphs[0], str(n), 9, color=C["muted"], align=PP_ALIGN.RIGHT, space_after=0)

    def title_zone(sl, text):
        rect(sl, 0, 0, Inches(0.16), SH, C["accent"])            # brand edge stripe
        tf = tb(sl, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.9))
        para(tf.paragraphs[0], text, 26, font=FH, bold=True, color=C["accent"])
        rect(sl, Inches(0.72), Inches(1.28), Inches(2.2), Pt(2.2), C["accent2"])  # underline

    def notes(sl, text):
        if text:
            sl.notes_slide.notes_text_frame.text = str(text)

    def flow_box(sl, l, t, w, h, text):
        sp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        sp.shadow.inherit = False
        sp.fill.solid(); sp.fill.fore_color.rgb = C["band"]
        sp.line.color.rgb = C["accent"]; sp.line.width = Pt(1.5)
        tf = sp.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(6); tf.margin_top = tf.margin_bottom = Pt(3)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf.paragraphs[0], str(text), 12, font=FB, bold=False, color=C["ink"],
             align=PP_ALIGN.CENTER, space_after=0)
        return sp

    def flow_arrow(sl, l, t, w, h, down=False):
        shp = MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW
        a = sl.shapes.add_shape(shp, l, t, w, h)
        a.shadow.inherit = False
        a.fill.solid(); a.fill.fore_color.rgb = C["accent2"]
        a.line.fill.background()
        return a

    def card_panel(sl, l, t, w, h, heading, text, big=False):
        """An accent-bar card (subtle-border surface + top accent bar) — editable shapes."""
        panel = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        panel.shadow.inherit = False
        panel.fill.solid(); panel.fill.fore_color.rgb = C["band"]
        panel.line.color.rgb = C["surface"]; panel.line.width = Pt(1)
        rect(sl, l + Inches(0.12), t + Inches(0.12), w - Inches(0.24), Pt(3), C["accent2"])
        tf = tb(sl, l + Inches(0.28), t + Inches(0.32), w - Inches(0.56), h - Inches(0.5))
        para(tf.paragraphs[0], str(heading), 30 if big else 22, font=FH, bold=True,
             color=C["accent"], space_after=6)
        if big:
            track(tf.paragraphs[0], -0.5)
        if text:
            para(tf.add_paragraph(), str(text), 13, color=C["ink"], space_after=0)

    def flow_row(sl, steps, x0, y, total_w, box_h):
        """Lay a row of movable step boxes with arrows between them."""
        n = len(steps)
        arrow_w = Inches(0.5)
        gaps = (n - 1) * arrow_w
        box_w = int((total_w - gaps) / max(n, 1))
        x = x0
        for j, step in enumerate(steps):
            flow_box(sl, x, y, box_w, box_h, step)
            x += box_w
            if j < n - 1:
                flow_arrow(sl, x + Inches(0.05), y + int(box_h / 2) - Inches(0.14),
                           arrow_w - Inches(0.1), Inches(0.28))
                x += arrow_w

    def add_bullets(tf, items, size=18, ref=False):
        first = True
        for it in items:
            if isinstance(it, list):
                for sub in it:
                    p = tf.add_paragraph(); p.level = 1
                    r1 = p.add_run(); r1.text = "– "; r1.font.color.rgb = C["accent2"]
                    r2 = p.add_run(); r2.text = str(sub); r2.font.color.rgb = C["muted"]
                    for r in (r1, r2):
                        r.font.size = Pt(size - 3); r.font.name = FB
                    p.space_after = Pt(4)
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if ref:
                para(p, str(it), size, color=C["ink"], space_after=8)
            else:
                r1 = p.add_run(); r1.text = "▪  "; r1.font.color.rgb = C["accent"]
                r2 = p.add_run(); r2.text = str(it); r2.font.color.rgb = C["ink"]
                for r in (r1, r2):
                    r.font.size = Pt(size); r.font.name = FB
                p.space_after = Pt(9)
            first = False

    for i, s in enumerate(spec["slides"], 1):
        lay = s["layout"]
        sl = slide()

        if lay == "title":
            bg(sl)
            rect(sl, 0, 0, SW, Inches(0.28), C["accent"])          # top brand bar
            rect(sl, Inches(10.9), Inches(0.28), Inches(2.43), Inches(6.62), C["surface"])  # side panel
            rect(sl, Inches(10.9), Inches(0.28), Pt(3), Inches(6.62), C["accent2"])          # panel edge
            rect(sl, 0, Inches(6.9), SW, Inches(0.6), C["accent"])   # bottom band
            tf = tb(sl, Inches(0.9), Inches(2.1), Inches(9.7), Inches(2.9))
            para(tf.paragraphs[0], s.get("title", spec.get("title", "")), 46,
                 font=FH, bold=True, color=C["accent"], space_after=12)
            track(tf.paragraphs[0], -0.5); line_spacing(tf.paragraphs[0], 1.02)
            rect(sl, Inches(0.95), Inches(3.65), Inches(3.2), Pt(3.5), C["accent2"])
            for txt, sz, col in ((s.get("subtitle") or spec.get("subtitle"), 20, C["ink"]),
                                 (s.get("authors") or spec.get("authors"), 16, C["muted"]),
                                 (s.get("date") or spec.get("date"), 13, C["muted"])):
                if txt:
                    para(tf.add_paragraph(), txt, sz, color=col, space_after=4)

        elif lay == "section":
            bg(sl, C["accent"])
            rect(sl, 0, Inches(2.95), Inches(0.55), Inches(1.6), C["accent2"])  # side accent block
            tf = tb(sl, Inches(0.9), 0, Inches(11.5), SH, anchor=MSO_ANCHOR.MIDDLE)
            para(tf.paragraphs[0], s.get("title", ""), 40, font=FH, bold=True,
                 color=C["onaccent"])
            track(tf.paragraphs[0], -0.5)
            rect(sl, Inches(0.95), Inches(4.65), Inches(2.6), Pt(3.5), C["accent2"])
            pn = tb(sl, Inches(11.9), Inches(6.7), Inches(0.7), Inches(0.4))
            para(pn.paragraphs[0], str(i), 11, color=C["onaccent"], align=PP_ALIGN.RIGHT)

        elif lay in ("bullets", "references"):
            bg(sl); title_zone(sl, s.get("title", "References" if lay == "references" else ""))
            tf = tb(sl, Inches(0.8), Inches(1.6), Inches(11.8), Inches(5.2))
            add_bullets(tf, s.get("bullets") if lay == "bullets" else s.get("items", []),
                        size=15 if lay == "references" else 18, ref=(lay == "references"))
            footer(sl, i)

        elif lay == "two_column":
            bg(sl); title_zone(sl, s.get("title", ""))
            for col, x in (("left", Inches(0.8)), ("right", Inches(7.0))):
                rect(sl, x, Inches(1.62), Inches(5.55), Inches(5.1), C["band"])
                tf = tb(sl, x + Inches(0.25), Inches(1.8), Inches(5.05), Inches(4.8))
                add_bullets(tf, s.get(col, []), size=17)
            footer(sl, i)

        elif lay == "figure":
            bg(sl); title_zone(sl, s.get("title", ""))
            img = s.get("image")
            if img and os.path.exists(img):
                pic = sl.shapes.add_picture(img, Inches(1.4), Inches(1.7), height=Inches(4.5))
                if pic.width > Inches(10.5):
                    pic.width, pic.height = Inches(10.5), int(pic.height * Inches(10.5) / pic.width)
            else:
                para(tb(sl, Inches(1.4), Inches(3.2), Inches(10), Inches(1)).paragraphs[0],
                     f"[figure not found: {img}]", 16, color=C["muted"])
            if s.get("caption"):
                para(tb(sl, Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.6)).paragraphs[0],
                     s["caption"], 12, color=C["muted"])
            footer(sl, i)

        elif lay == "table":
            bg(sl); title_zone(sl, s.get("title", ""))
            data = s["table"]
            rows, cols = len(data), max(len(r) for r in data)
            gt = sl.shapes.add_table(rows, cols, Inches(0.8), Inches(1.7),
                                     Inches(11.7), Inches(min(0.5 * rows, 5.0))).table
            gt.first_row = True
            for r, row in enumerate(data):
                for c in range(cols):
                    cell = gt.cell(r, c)
                    cell.text = str(row[c]) if c < len(row) else ""
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (C["accent"] if r == 0
                                                else (C["band"] if r % 2 else C["bg"]))
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(12); run.font.name = FB
                            run.font.bold = (r == 0)
                            run.font.color.rgb = C["onaccent"] if r == 0 else C["ink"]
            footer(sl, i)

        elif lay == "flow":
            # BioRender-style experimental-design flow as NATIVE, MOVABLE shapes:
            # each step box and each arrow is its own editable object.
            bg(sl); title_zone(sl, s.get("title", "Experimental design"))
            lanes = s.get("lanes")
            if not lanes:
                lanes = [{"label": "", "steps": s.get("steps", [])}]
            top, bottom = Inches(1.7), Inches(6.8)
            avail_h = bottom - top
            lane_h = int(avail_h / max(len(lanes), 1))
            box_h = min(Inches(1.1), lane_h - Inches(0.5))
            for li, lane in enumerate(lanes):
                ly = top + li * lane_h
                label = lane.get("label", "")
                if label:
                    ltf = tb(sl, Inches(0.7), ly, Inches(1.7), lane_h, anchor=MSO_ANCHOR.MIDDLE)
                    para(ltf.paragraphs[0], label, 12, font=FH, bold=True, color=C["accent"],
                         space_after=0)
                    x0, row_w = Inches(2.5), Inches(10.1)
                else:
                    x0, row_w = Inches(0.8), Inches(11.7)
                flow_row(sl, lane.get("steps", []), x0,
                         ly + int((lane_h - box_h) / 2), row_w, box_h)
                if li < len(lanes) - 1:  # arrow down to next lane
                    flow_arrow(sl, Inches(6.5), ly + lane_h - Inches(0.32),
                               Inches(0.34), Inches(0.3), down=True)
            footer(sl, i)

        elif lay == "result":
            # Results slide: figure OR table on the left, key findings on the right.
            bg(sl); title_zone(sl, s.get("title", "Results"))
            img, data = s.get("image"), s.get("table")
            if data:
                rows, cols = len(data), max(len(r) for r in data)
                gt = sl.shapes.add_table(rows, cols, Inches(0.8), Inches(1.75),
                                         Inches(7.4), Inches(min(0.5 * rows, 4.6))).table
                gt.first_row = True
                for r, row in enumerate(data):
                    for c in range(cols):
                        cell = gt.cell(r, c)
                        cell.text = str(row[c]) if c < len(row) else ""
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = (C["accent"] if r == 0
                                                    else (C["band"] if r % 2 else C["bg"]))
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                run.font.size = Pt(11); run.font.name = FB
                                run.font.bold = (r == 0)
                                run.font.color.rgb = C["onaccent"] if r == 0 else C["ink"]
            elif img and os.path.exists(img):
                pic = sl.shapes.add_picture(img, Inches(0.8), Inches(1.75), width=Inches(7.4))
                if pic.height > Inches(4.9):
                    pic.height, pic.width = Inches(4.9), int(pic.width * Inches(4.9) / pic.height)
            else:
                para(tb(sl, Inches(0.8), Inches(3.2), Inches(7.2), Inches(1)).paragraphs[0],
                     f"[figure not found: {img}]", 14, color=C["muted"])
            if s.get("caption"):
                para(tb(sl, Inches(0.8), Inches(6.5), Inches(7.4), Inches(0.5)).paragraphs[0],
                     s["caption"], 11, color=C["muted"])
            # Key findings panel (right)
            rect(sl, Inches(8.5), Inches(1.75), Inches(4.1), Inches(4.95), C["band"])
            rect(sl, Inches(8.5), Inches(1.75), Inches(4.1), Inches(0.5), C["accent"])
            kh = tb(sl, Inches(8.7), Inches(1.8), Inches(3.8), Inches(0.42),
                    anchor=MSO_ANCHOR.MIDDLE)
            para(kh.paragraphs[0], "Key findings", 14, font=FH, bold=True,
                 color=C["onaccent"], space_after=0)
            kf = tb(sl, Inches(8.7), Inches(2.35), Inches(3.75), Inches(4.2))
            add_bullets(kf, s.get("key_findings", []), size=13)
            footer(sl, i)

        elif lay == "metric":
            # Big-number hero: one oversized figure the audience remembers.
            bg(sl, C["surface"])
            if s.get("title"):
                title_zone(sl, s["title"])
            vtf = tb(sl, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.4),
                     anchor=MSO_ANCHOR.MIDDLE)
            para(vtf.paragraphs[0], str(s.get("value", "")), 96, font=FH, bold=True,
                 color=C["accent"], align=PP_ALIGN.CENTER, space_after=4)
            track(vtf.paragraphs[0], -1.0)
            if s.get("label"):
                para(vtf.add_paragraph(), str(s["label"]), 24, color=C["ink"],
                     align=PP_ALIGN.CENTER, space_after=2)
            if s.get("support"):
                para(vtf.add_paragraph(), str(s["support"]), 15, color=C["muted"],
                     align=PP_ALIGN.CENTER, space_after=0)
            footer(sl, i)

        elif lay == "cards":
            # Metrics/feature grid of accent-bar cards (up to 6).
            bg(sl); title_zone(sl, s.get("title", ""))
            cards = s.get("cards", [])[:6]
            n = len(cards)
            cols = 1 if n <= 1 else (2 if n in (2, 4) else 3)
            rows_n = max(1, math.ceil(n / cols))
            gap = Inches(0.35)
            x0, y0, area_w, area_h = Inches(0.8), Inches(1.85), Inches(11.7), Inches(4.85)
            cw = int((area_w - gap * (cols - 1)) / cols)
            ch = int((area_h - gap * (rows_n - 1)) / rows_n)
            big = rows_n == 1  # single row of headline stats → oversized headings
            for idx, card in enumerate(cards):
                r, c = divmod(idx, cols)
                cx = x0 + c * (cw + gap)
                cy = y0 + r * (ch + gap)
                card_panel(sl, cx, cy, cw, ch, card.get("heading", ""),
                           card.get("text", ""), big=big)
            footer(sl, i)

        notes(sl, s.get("notes"))

    prs.save(out_path)
    return out_path


def selftest():
    spec = {"title": "T", "theme": "sage", "slides": [
        {"layout": "title", "title": "T", "subtitle": "S"},
        {"layout": "section", "title": "Part 1"},
        {"layout": "bullets", "title": "Points", "bullets": ["one", ["sub"], "two"], "notes": "n"},
        {"layout": "two_column", "title": "C", "left": ["a"], "right": ["b"]},
        {"layout": "table", "title": "Data", "table": [["Trt", "Value"], ["A", "1.2"]]},
        {"layout": "flow", "title": "Design", "steps": ["Sample", "Assay", "Analyse"]},
        {"layout": "result", "title": "Result 1", "table": [["Trt", "Val"], ["A", "1.2"]],
         "key_findings": ["A beat control", "p < 0.05"]},
        {"layout": "metric", "title": "Headline", "value": "42%",
         "label": "reduction in spoilage", "support": "vs control (p < 0.05)"},
        {"layout": "cards", "title": "At a glance", "cards": [
            {"heading": "n = 1,240", "text": "participants"},
            {"heading": "6 trials", "text": "pooled"},
            {"heading": "2019-2024", "text": "window"}]},
        {"layout": "references", "title": "References", "items": ["Ref 1"]}]}
    assert validate_spec(spec) == [], validate_spec(spec)
    bad = validate_spec({"slides": [{"layout": "bogus"}, {"layout": "figure"},
                                    {"layout": "metric"}, {"layout": "cards"}]})
    assert any("layout 'bogus'" in p for p in bad) and any("missing 'image'" in p for p in bad)
    assert any("needs a 'value'" in p for p in bad) and any("needs 'cards'" in p for p in bad)
    assert validate_spec({}) == ["spec.slides must be a non-empty list"]
    # themes.json loads and default resolves
    t = load_theme("sage"); assert t["accent"] and t["font_body"]
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("OK: build_pptx selftest passed (validator+theme; python-pptx not installed)")
        return
    import tempfile
    out = os.path.join(tempfile.mkdtemp(), "t.pptx")
    build(spec, out)
    with zipfile.ZipFile(out) as z:
        names = z.namelist()
        assert "ppt/presentation.xml" in names
        s3 = z.read("ppt/slides/slide3.xml").decode("utf-8", "ignore")
        assert "<a:t>" in s3 and "one" in s3, "bullets must be native editable runs"
        s5 = z.read("ppt/slides/slide5.xml").decode("utf-8", "ignore")
        assert "<a:tbl>" in s5, "data slide must be a native table"
        s8 = z.read("ppt/slides/slide8.xml").decode("utf-8", "ignore")
        assert "42%" in s8 and "spc=" in s8, "metric slide: editable value with tracking"
        s9 = z.read("ppt/slides/slide9.xml").decode("utf-8", "ignore")
        assert s9.count("roundRect") >= 3 and "n = 1,240" in s9, "cards: native card shapes"
    print("OK: build_pptx selftest passed (validator + themed editable .pptx built)")


def main(argv):
    if "--selftest" in argv:
        selftest(); return 0
    args = [a for a in argv[1:] if not a.startswith("--")]
    out = theme = None
    for i, a in enumerate(argv):
        if a == "--out": out = argv[i + 1]
        if a == "--theme": theme = argv[i + 1]
    if not args or not out:
        print("usage: build_pptx.py <deck.json> --out <deck.pptx> [--theme <name>] | --selftest")
        return 1
    spec = json.load(open(args[0]))
    problems = validate_spec(spec)
    if problems:
        print("FAIL: invalid deck spec:")
        for p in problems:
            print("  ✗", p)
        return 1
    path = build(spec, out, theme)
    print(f"OK: built editable {path} ({len(spec['slides'])} slides, theme "
          f"'{theme or spec.get('theme') or 'default'}')")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
